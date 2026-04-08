from pathlib import Path
from shutil import copy2

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parents[3]
FIG = BASE / 'figures'
SRC = Path('/Users/demn/Documents/autoru.pptx')
BACKUP = Path('/Users/demn/Documents/autoru_backup_before_rework.pptx')

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(25, 25, 25)
MUTED = RGBColor(95, 99, 104)
BORDER = RGBColor(166, 166, 166)
LIGHT = RGBColor(248, 248, 248)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def backup_file():
    if SRC.exists() and not BACKUP.exists():
        copy2(SRC, BACKUP)


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font_name='Arial'):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_header(slide, title, number=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.92))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()
    add_text(slide, Inches(0.42), Inches(0.08), Inches(12.0), Inches(0.52), title, font_size=24, bold=True, color=WHITE, font_name='Arial Narrow')
    if number is not None:
        add_text(slide, Inches(12.7), Inches(6.9), Inches(0.28), Inches(0.18), str(number), font_size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_panel(slide, left, top, width, height, title, lines, body_size=16):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Arial'
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = TEXT
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Arial'
        p.font.size = Pt(body_size)
        p.font.color.rgb = TEXT
        p.bullet = True
        p.level = 0
    return shape


def add_picture_fit(slide, path, left, top, width, height):
    path = str(path)
    with Image.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    box_aspect = width / height
    if aspect > box_aspect:
        pic_w = width
        pic_h = width / aspect
        pic_left = left
        pic_top = top + (height - pic_h) / 2
    else:
        pic_h = height
        pic_w = height * aspect
        pic_top = top
        pic_left = left + (width - pic_w) / 2
    return slide.shapes.add_picture(path, pic_left, pic_top, width=pic_w, height=pic_h)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)
    add_text(slide, Inches(0.35), Inches(0.12), Inches(3.0), Inches(0.3), 'Студенческая работа', font_size=16, color=WHITE)
    add_text(slide, Inches(1.15), Inches(2.0), Inches(11.0), Inches(1.3), 'Факторы цены автомобилей\nна Auto.ru', font_size=31, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.35), Inches(4.25), Inches(6.7), Inches(0.7), 'АНДАН: данные, гипотезы и критерии', font_size=21, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4.55), Inches(5.82), Inches(4.2), Inches(0.35), 'Никита Дёмин', font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    return slide


def section_slide(prs, subtitle, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(3.75), SLIDE_H)
    band.fill.solid()
    band.fill.fore_color.rgb = BLACK
    band.line.fill.background()
    add_text(slide, Inches(4.9), Inches(2.55), Inches(4.0), Inches(0.5), subtitle, font_size=22, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.75), Inches(3.35), Inches(6.0), Inches(0.9), title, font_size=35, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    return slide


def graph_slide(prs, number, title, figure_name, note_title, note_lines, fig_left=0.85, fig_top=1.28, fig_w=11.6, fig_h=4.42):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, title, number)
    add_picture_fit(slide, FIG / figure_name, Inches(fig_left), Inches(fig_top), Inches(fig_w), Inches(fig_h))
    add_panel(slide, Inches(0.85), Inches(5.92), Inches(11.45), Inches(0.7), note_title, note_lines, body_size=15)
    return slide


def build():
    backup_file()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Цель и задачи', 2)
    add_panel(slide, Inches(0.75), Inches(1.28), Inches(12.0), Inches(1.25), 'Цель', [
        'Построить исследовательский каркас для анализа рынка Auto.ru: проверить форму распределения цены, неоднородность сегментов, структуру факторов и подготовить корректный переход к регрессионному моделированию.'
    ], body_size=17)
    add_panel(slide, Inches(0.75), Inches(2.90), Inches(12.0), Inches(3.20), 'Задачи', [
        'Проверить корректность и ограничения датасета до начала анализа.',
        'Сформулировать последовательные гипотезы о распределении, сегментации и структуре рынка.',
        'Выбрать статистические критерии, согласованные с видом распределений.',
        'Выявить ключевые и proxy-признаки, чтобы отделить реальные факторы цены от технических дубликатов.',
        'Сформировать постановку следующего этапа: сравнение baseline-регрессоров на raw и processed данных.'
    ], body_size=16)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Обзор данных', 3)
    add_panel(slide, Inches(0.75), Inches(1.25), Inches(5.9), Inches(5.05), 'Что находится в датасете', [
        'После очистки и приведения типов в анализе осталось 6201 объявление.',
        'В cleaned-датасете 27 полей: исходные характеристики карточки + производные признаки очистки.',
        'Целевая переменная исследования: price.',
        'Для распределений и факторов анализируются как исходная цена, так и преобразование price_log1p.'
    ], body_size=16)
    add_panel(slide, Inches(6.75), Inches(1.25), Inches(5.8), Inches(5.05), 'Смысловые группы признаков', [
        'Структурные: brand, model, generation, body_type, condition.',
        'Технические: engine_power_hp, engine_volume, transmission, fuel_type, drive_type.',
        'Эксплуатационные: age, mileage.',
        'Служебные и технические поля очистки: description_text, url, parsed_at, commercial_* и др.'
    ], body_size=16)

    graph_slide(
        prs, 4, 'Проверка корректности и ограничения датасета', 'presentation_dataset_quality.png', 'Вывод',
        ['Ключевые ограничения задаются не размером выборки, а пропусками и кардинальностью: owners_count пустой, engine_volume и mileage частично пропущены, а brand/model/generation очень детализированы.']
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Гипотезы 1–3: распределение и сегментация', 5)
    add_panel(slide, Inches(0.72), Inches(1.20), Inches(12.0), Inches(1.50), 'H1. Цена не образует один нормальный закон', [
        'Формулировка: распределение price ближе к смеси правосторонне-асимметричных подраспределений, чем к единой нормали.',
        'Почему важно: от этого зависит выбор критериев и устойчивость будущей функции потерь.',
        'Как проверяем: histogram + KDE, Q-Q plot, Shapiro–Wilk, D’Agostino, Anderson–Darling.'
    ], body_size=14)
    add_panel(slide, Inches(0.72), Inches(3.02), Inches(12.0), Inches(1.50), 'H2. Логарифмирование улучшает форму, но не устраняет рыночную смесь', [
        'Формулировка: log(price) уменьшает асимметрию и тяжесть хвоста, но не превращает рынок в идеальную нормаль.',
        'Почему важно: это влияет на постановку регрессии в исходной шкале и в log-space.',
        'Как проверяем: сравнение histogram/Q-Q/raw vs log, сопоставление skewness и kurtosis.'
    ], body_size=14)
    add_panel(slide, Inches(0.72), Inches(4.84), Inches(12.0), Inches(1.50), 'H3. Рынок неоднороден уже на уровне new/used', [
        'Формулировка: объединение новых и автомобилей с пробегом в один рынок искажает форму распределения и силу факторов.',
        'Почему важно: pooled-анализ скрывает разные механизмы ценообразования.',
        'Как проверяем: density overlay, Mann–Whitney U, сегментные корреляции.'
    ], body_size=14)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Гипотезы 4–6: факторная структура рынка', 6)
    add_panel(slide, Inches(0.72), Inches(1.20), Inches(12.0), Inches(1.50), 'H4. Эффект износа нелинеен и сегментно-зависим', [
        'Формулировка: дисконт по age и mileage формируется неравномерно, причём в used-сегменте он сильнее, чем в new.',
        'Почему важно: линейные коэффициенты могут недоописывать реальную форму связи.',
        'Как проверяем: scatter + LOWESS, Spearman по сегментам new и used.'
    ], body_size=14)
    add_panel(slide, Inches(0.72), Inches(3.02), Inches(12.0), Inches(1.50), 'H5. Структурные признаки сильнее отдельных технических', [
        'Формулировка: brand/model/generation/condition объясняют цену лучше, чем isolated технические признаки.',
        'Почему важно: рынок позиционируется комбинацией класса и состояния, а не одной колонкой.',
        'Как проверяем: Cramér’s V, chi-square, Spearman, bar chart силы связи.'
    ], body_size=14)
    add_panel(slide, Inches(0.72), Inches(4.84), Inches(12.0), Inches(1.50), 'H6. Часть признаков является proxy-переменными', [
        'Формулировка: некоторые поля отражают уже учтённую рыночную структуру, а не добавляют новый сигнал.',
        'Почему важно: proxy-признаки мешают интерпретации и дублируют информацию в модели.',
        'Как проверяем: Cramér’s V matrix, корреляционная матрица, VIF.'
    ], body_size=14)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Гипотезы 7–8: подготовка к моделированию', 7)
    add_panel(slide, Inches(0.72), Inches(1.30), Inches(12.0), Inches(1.65), 'H7. Preprocessing должен стабилизировать статистические выводы', [
        'Формулировка: логарифмирование цены и очистка выбросов должны уменьшать влияние хвоста и делать baseline-оценки устойчивее.',
        'Почему важно: если preprocessing не улучшает устойчивость, он не оправдан как исследовательский этап.',
        'Как проверяем дальше: одинаковые split и baseline-модели на raw и processed данных.'
    ], body_size=14)
    add_panel(slide, Inches(0.72), Inches(3.25), Inches(12.0), Inches(1.65), 'H8. Единой линейной спецификации для такого рынка недостаточно', [
        'Формулировка: тяжёлые хвосты, сегментация new/used и нелинейный эффект износа должны ограничивать чисто линейную модель.',
        'Почему важно: выбор baseline-моделей и функции потерь должен соответствовать структуре рынка, а не только удобству интерпретации.',
        'Как проверяем дальше: сравнение линейных и более гибких регрессоров в одной исследовательской схеме.'
    ], body_size=14)
    add_panel(slide, Inches(0.72), Inches(5.20), Inches(12.0), Inches(1.00), 'Исследовательская цепочка', [
        'Корректность данных -> распределения -> критерии -> однородность -> факторы -> proxy-признаки -> постановка регрессии.'
    ], body_size=15)

    section_slide(prs, 'Том 1', 'Распределения и критерии')

    graph_slide(
        prs, 9, 'Структура выборки по сегментам new/used', 'condition_counts_presentation_ru.png', 'Вывод',
        ['Обе группы достаточно крупные для статистических проверок. Это позволяет анализировать неоднородность new и used не как визуальное впечатление, а как формальную гипотезу.']
    )

    graph_slide(
        prs, 10, 'Структура выборки по ценовым сегментам', 'presentation_price_bands.png', 'Вывод',
        ['По ценовому профилю выборка распадается как минимум на массовый, средний и дорогой подрынки. Уже здесь цена выглядит как смесь сегментов, а не как один общий режим.']
    )

    graph_slide(
        prs, 11, 'Распределение price и log(price)', 'presentation_price_vs_log_distribution.png', 'Вывод',
        ['В исходной шкале цена имеет выраженную правостороннюю асимметрию и тяжёлый премиальный хвост. После логарифмирования центральная масса становится компактнее, но рынок всё равно не превращается в идеальную нормаль.']
    )

    graph_slide(
        prs, 12, 'Q-Q plot: price и log(price)', 'presentation_qq_raw_vs_log.png', 'Вывод',
        ['В исходной шкале хвосты резко уходят от прямой. В лог-шкале центральная часть выравнивается, но крайние квантили всё ещё отклоняются, что соответствует log-normal-like смеси, а не чистому гауссовскому закону.']
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Как интерпретировалась нормальность', 13)
    add_panel(slide, Inches(0.75), Inches(1.25), Inches(5.9), Inches(5.10), 'Формальная проверка', [
        'H0: распределение цены согласуется с нормальным законом.',
        'price: Shapiro–Wilk W = 0.71, D’Agostino K² = 4989.91, Anderson A² = 335.29 -> H0 отвергается.',
        'log(price): W = 0.97, K² = 275.78, A² = 83.83 -> форма заметно улучшается, но H0 всё равно отвергается.',
        'При N > 5000 одного p-value недостаточно, поэтому решение принималось по совокупности: hist + Q-Q + skewness + kurtosis.'
    ], body_size=15)
    add_panel(slide, Inches(6.75), Inches(1.25), Inches(5.8), Inches(5.10), 'Содержательная интерпретация', [
        'price: skew = 3.68, kurtosis = 25.96 -> сильная правая асимметрия и тяжёлый хвост.',
        'log(price): skew = -0.50, kurtosis = 3.53 -> хвост сглажен, но идеальной гауссовой формы нет.',
        'Для автомобильного рынка это естественно: масса объявлений сосредоточена в массовом сегменте, а премиум формирует длинный правый хвост.',
        'Практический вывод: для групповых сравнений нужны непараметрические критерии, а для моделирования разумно рассматривать log-space и устойчивые loss functions.'
    ], body_size=15)

    graph_slide(
        prs, 14, 'Сравнение распределений цены в new и used', 'presentation_homogeneity_new_used_overlay.png', 'Вывод',
        ['Распределения различаются не только по медиане, но и по форме: new-сегмент сдвинут вправо, а used-сегмент заметно шире и сильнее растянут в дешёвый диапазон.']
    )

    graph_slide(
        prs, 15, 'Цена по condition: форма распределения внутри групп', 'presentation_condition_violin.png', 'Вывод',
        ['У used-сегмента распределение существенно шире и менее концентрировано. Это уже не просто сдвиг центра, а другой режим ценообразования с большей дисперсией и большим влиянием износа.']
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Как проверялась однородность', 16)
    add_panel(slide, Inches(0.75), Inches(1.25), Inches(5.9), Inches(5.10), 'Критерий и гипотезы', [
        'H0: распределения цены в new и used совпадают.',
        'H1: распределения различаются.',
        'Так как цена ненормальна и асимметрична, для сравнения использовался Mann–Whitney U, а не t-test по средним.',
        'Критерий сравнивает ранговые распределения и не требует нормальности выборок.'
    ], body_size=15)
    add_panel(slide, Inches(6.75), Inches(1.25), Inches(5.8), Inches(5.10), 'Результат и смысл', [
        'U = 8 277 249, p-value < 0.001 -> H0 отвергается.',
        'Медиана new = 3.50 млн руб., медиана used = 1.45 млн руб.',
        'Следовательно, единый pooled-анализ по всем авто неадекватен: далее рынок нужно интерпретировать как минимум с учётом разделения на new и used.',
        'Это решение напрямую влияет на следующую постановку регрессии.'
    ], body_size=15)

    graph_slide(
        prs, 17, 'Цена по body_type', 'presentation_body_type_boxplot.png', 'Вывод',
        ['Кузов отражает не только форму автомобиля, но и рыночный класс: внедорожники, пикапы и купе смещены выше седанов и хэтчбеков. Значит, body_type участвует в структурном позиционировании цены.']
    )

    graph_slide(
        prs, 18, 'Цена по top-N брендам', 'presentation_brand_boxplot.png', 'Вывод',
        ['Бренд уже сам по себе задаёт разные уровни медианы и разную ширину распределения. Это аргумент в пользу того, что структурные идентификаторы рынка сильнее отдельных технических полей.']
    )

    graph_slide(
        prs, 19, 'Цена по transmission, fuel_type и drive_type', 'presentation_small_category_boxplots.png', 'Вывод',
        ['Технические категории тоже разделяют рынок, но слабее структурных. Например, автомат и гибрид/электро чаще связаны с дорогим сегментом, однако это скорее отражение позиционирования модели, чем самостоятельная первопричина цены.']
    )

    section_slide(prs, 'Том 2', 'Факторы и зависимости')

    graph_slide(
        prs, 21, 'Числовые признаки: сила связи с ценой', 'presentation_numeric_associations.png', 'Вывод',
        ['Главный числовой сигнал идёт от engine_power_hp, возраста и пробега. Это не просто наличие связи, а устойчивый рыночный градиент: мощность повышает цену, а износ и mileage её понижают.']
    )

    graph_slide(
        prs, 22, 'Числовые факторы внутри new и used', 'presentation_segmented_effects.png', 'Вывод',
        ['Внутри сегментов действуют разные силы: в new-маркете мощность особенно сильна, а в used доминируют age и mileage. Это прямое подтверждение гипотезы о разных режимах ценообразования.']
    )

    graph_slide(
        prs, 23, 'Возраст: исходная и логарифмическая шкала', 'presentation_dep_age_dual.png', 'Вывод',
        ['Связь с возрастом отрицательная и нелинейная: основной дисконт формируется на ранних годах эксплуатации, после чего темп снижения цены замедляется.']
    )

    graph_slide(
        prs, 24, 'Пробег: исходная и логарифмическая шкала', 'presentation_dep_mileage_dual.png', 'Вывод',
        ['Пробег даёт убывающий, но не линейный эффект: основной дисконт приходится на первые диапазоны роста mileage, затем кривая становится более пологой.']
    )

    graph_slide(
        prs, 25, 'Мощность двигателя: исходная и логарифмическая шкала', 'presentation_dep_power_dual.png', 'Вывод',
        ['Мощность связана с ценой положительно и монотонно, но не строго линейно: в верхнем диапазоне прирост цены становится сегментно-зависимым и отражает переход к премиальному классу.']
    )

    graph_slide(
        prs, 26, 'Объём двигателя: исходная и логарифмическая шкала', 'presentation_dep_volume_dual.png', 'Вывод',
        ['Объём двигателя объясняет цену слабее мощности и возраста. Это вспомогательный технический фактор, а не главный структурный драйвер рынка.']
    )

    graph_slide(
        prs, 27, 'Категориальные признаки: сила связи с ценой', 'presentation_categorical_associations.png', 'Вывод',
        ['По Cramér’s V лидируют model, generation, condition и brand. То есть рыночное позиционирование автомобиля объясняет цену сильнее, чем отдельная техническая характеристика или декоративное поле.']
    )

    graph_slide(
        prs, 28, 'Тепловая карта числовых зависимостей', 'presentation_numeric_corr_heatmap.png', 'Вывод',
        ['Матрица показывает две ключевые вещи: price_log1p связан с age/mileage/engine_power_hp, а year и age почти зеркальны. Это сигнал не только о факторах цены, но и о дублировании признаков.']
    )

    graph_slide(
        prs, 29, 'Матрица Cramér’s V для категориальных признаков', 'presentation_categorical_cramers_heatmap.png', 'Вывод',
        ['Самый жёсткий proxy-эффект здесь — condition и seller_type. Высокие связи condition с body_type и transmission тоже отражают сегментную структуру, а не независимые причинные эффекты.']
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Proxy-признаки и дублирование информации', 30)
    add_panel(slide, Inches(0.75), Inches(1.25), Inches(5.85), Inches(5.10), 'Что оказалось избыточным', [
        'seller_type почти полностью дублирует condition: Cramér’s V = 1.00.',
        'year и age кодируют одну и ту же информацию в противоположных направлениях: Spearman = -1.00.',
        'owners_count не включается в анализ как фактор цены, потому что поле полностью пустое.'
    ], body_size=15)
    add_panel(slide, Inches(6.70), Inches(1.25), Inches(5.85), Inches(5.10), 'Почему это важно', [
        'Proxy-поля создают ложное ощущение большого числа факторов, хотя реально рынок описывается меньшим числом независимых осей.',
        'В статистических выводах это мешает интерпретации, а в моделировании может приводить к нестабильным коэффициентам и избыточной размерности.',
        'Поэтому до регрессии часть признаков должна быть исключена или объединена концептуально.'
    ], body_size=15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Итоговый аналитический набор признаков', 31)
    add_panel(slide, Inches(0.75), Inches(1.25), Inches(3.9), Inches(5.15), 'Структурное ядро', [
        'brand',
        'model',
        'generation',
        'condition',
        'body_type'
    ], body_size=17)
    add_panel(slide, Inches(4.55), Inches(1.25), Inches(3.8), Inches(5.15), 'Технический слой', [
        'engine_power_hp',
        'engine_volume',
        'transmission',
        'fuel_type',
        'drive_type'
    ], body_size=17)
    add_panel(slide, Inches(8.35), Inches(1.25), Inches(4.2), Inches(5.15), 'Износ и исключения', [
        'age',
        'mileage',
        'Исключены: seller_type, year, owners_count, color, region, customs, pts_type, steering_wheel, commercial_* и служебные поля.'
    ], body_size=15)

    section_slide(prs, 'Том 3', 'Переход к моделированию')

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Почему следующий этап ставится как регрессия', 33)
    add_panel(slide, Inches(0.75), Inches(1.25), Inches(3.95), Inches(5.10), 'Постановка задачи', [
        'Целевая переменная непрерывна, поэтому задача естественно формулируется как регрессия.',
        'При этом регрессия не равна одной линейной модели: на таком рынке важно сравнивать разные семейства baseline-подходов.',
        'Рассматриваться будут и исходная шкала price, и логарифмическая шкала log(1 + price).'
    ], body_size=15)
    add_panel(slide, Inches(4.60), Inches(1.25), Inches(3.95), Inches(5.10), 'Какие baseline-модели разумны', [
        'Linear Regression / Ridge / Lasso как интерпретируемый базовый уровень.',
        'Random Forest Regressor как непараметрический ансамбль, чувствительный к взаимодействиям признаков.',
        'Gradient Boosting / XGBoost как более гибкие модели для нелинейностей и иерархических эффектов.'
    ], body_size=15)
    add_panel(slide, Inches(8.45), Inches(1.25), Inches(4.10), Inches(5.10), 'Почему это вытекает из анализа', [
        'Мы уже увидели heavy tail, сегментацию new/used и нелинейный эффект износа.',
        'Следовательно, единая линейная спецификация заранее не должна считаться достаточной.',
        'Это не обещание лучшего качества, а исследовательская гипотеза следующего этапа.'
    ], body_size=15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Функция потерь и выбор шкалы', 34)
    add_panel(slide, Inches(0.75), Inches(1.25), Inches(3.85), Inches(5.10), 'Линейные baseline-модели', [
        'OLS минимизирует L = (1/n) Σ (yᵢ - ŷᵢ)².',
        'Ridge: L + λ||β||²₂.',
        'Lasso: L + λ||β||₁.',
        'Квадратичная ошибка чувствительна к дорогому хвосту и редким выбросам.'
    ], body_size=15)
    add_panel(slide, Inches(4.55), Inches(1.25), Inches(3.95), Inches(5.10), 'Гибкие регрессоры', [
        'Random Forest выбирает разбиения по уменьшению дисперсии / squared_error.',
        'Gradient Boosting и XGBoost итеративно минимизируют регрессионный loss на остатках.',
        'Для heavy-tail рынка особенно важна устойчивость к редким дорогим объявлениям.'
    ], body_size=15)
    add_panel(slide, Inches(8.45), Inches(1.25), Inches(4.10), Inches(5.10), 'Почему важен log-space', [
        'Если оптимизировать на z = log(1 + price), премиальный хвост перестаёт доминировать так сильно.',
        'MAE и Huber потенциально устойчивее к выбросам, чем чистый MSE.',
        'Следовательно, в следующем этапе нужно сравнивать не только модели, но и шкалы/потери.'
    ], body_size=15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Дизайн следующего этапа: raw vs processed', 35)
    add_panel(slide, Inches(0.75), Inches(1.25), Inches(3.95), Inches(5.10), 'Что сравнивается', [
        'Один и тот же train/test split для raw и processed вариантов датасета.',
        'Один и тот же набор baseline-регрессоров.',
        'Пул по всем авто и, при необходимости, сегментный вариант с учётом new/used.'
    ], body_size=15)
    add_panel(slide, Inches(4.60), Inches(1.25), Inches(3.95), Inches(5.10), 'Метрики качества', [
        'MAE и RMSE — абсолютная ошибка в исходной шкале.',
        'MAPE — относительная ошибка для практической интерпретации.',
        'R² — доля объяснённой вариации.',
        'Отдельно оценивается устойчивость выводов к хвостам и выбросам.'
    ], body_size=15)
    add_panel(slide, Inches(8.45), Inches(1.25), Inches(4.10), Inches(5.10), 'Какая исследовательская цель', [
        'Количественно оценить, улучшает ли preprocessing предсказуемость цены, а не просто улучшает графики.',
        'Проверить, где именно линейные модели ограничены: хвосты, сегментация, взаимодействия или пропуски.',
        'Только после этого делать выводы о целесообразности более гибкого ML-подхода.'
    ], body_size=15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, 'Выводы', 36)
    add_panel(slide, Inches(0.75), Inches(1.20), Inches(12.0), Inches(5.30), 'Основные результаты первого этапа', [
        'Цена на рынке Auto.ru не согласуется с единой нормальной моделью и ближе к log-normal-like смеси сегментных распределений.',
        'Логарифмирование заметно улучшает форму распределения, но не устраняет неоднородность рынка.',
        'Сегменты new и used статистически различаются по форме и уровню распределения цены, поэтому pooled-анализ ограничен.',
        'Структурные признаки brand/model/generation/condition объясняют рынок лучше, чем отдельные вторичные поля.',
        'Часть признаков оказалась proxy-переменными и была исключена до перехода к регрессии.',
        'Следующий этап проекта — формализованное сравнение baseline-регрессоров и функций потерь на raw и processed данных.'
    ], body_size=15)

    prs.save(str(SRC))
    print(SRC)
    print('slides', len(prs.slides))


if __name__ == '__main__':
    build()
